"""디자인 시스템 — 10 테마 + 36 LAYOUTS 검증."""

import pytest


# ════════════════════════════════════════════════════════════
# _to_emu — float 좌표 자동 변환 (회귀 방지)
# ════════════════════════════════════════════════════════════

class TestToEmuHelper:
    """``_to_emu()`` — raw float 좌표를 자동 EMU 변환.

    이전에 사용자가 ``T(s, 0.6, 4.0, 7.07, 1.5, ...)`` 처럼 인치 float 을
    전달했을 때 거의 0 EMU 로 해석되어 모든 컴포넌트가 슬라이드 상단에
    겹치는 깨짐이 있었음. 이 회귀 방지 테스트.
    """

    def test_float_converted_to_inches(self):
        from src.generators.slide_kit import _to_emu
        from pptx.util import Inches
        assert _to_emu(2.5) == Inches(2.5)
        assert _to_emu(0.6) == Inches(0.6)

    def test_int_under_1000_converted_to_inches(self):
        from src.generators.slide_kit import _to_emu
        from pptx.util import Inches
        assert _to_emu(2) == Inches(2)

    def test_emu_int_passthrough(self):
        """이미 EMU 인 큰 int 는 그대로."""
        from src.generators.slide_kit import _to_emu
        from pptx.util import Inches
        emu_value = Inches(2.5)  # 2286000
        assert _to_emu(emu_value) == emu_value

    def test_none_passthrough(self):
        from src.generators.slide_kit import _to_emu
        assert _to_emu(None) is None

    def test_inches_object_passthrough(self):
        from src.generators.slide_kit import _to_emu
        from pptx.util import Inches
        v = Inches(3.0)
        assert _to_emu(v) == v


class TestComponentsAcceptFloat:
    """KPIS/COLS/HIGHLIGHT/QUOTE/ICON_CARDS/METRIC_CARD 가 float 인치 좌표 수용."""

    def _new_slide(self, format_name="legacy_16_9"):
        from src.generators import slide_kit
        slide_kit.apply_format(format_name)
        prs = slide_kit.new_presentation()
        return slide_kit, prs, slide_kit.new_slide(prs)

    def test_kpis_with_float_y(self):
        sk, _, s = self._new_slide()
        sk.KPIS(s, items=[{"value": "30%", "label": "X"}], y=2.5)
        # 마지막 추가된 shape 의 y 검증
        rounded = [sh for sh in s.shapes if "Rounded" in sh.name]
        assert rounded
        y_in = rounded[0].top / 914400
        assert y_in > 2.0, f"KPIS y 변환 실패: y_in={y_in}"

    def test_metric_card_with_float_coords(self):
        sk, _, s = self._new_slide()
        sk.METRIC_CARD(s, 0.6, 2.5, 7.07, 3.5, value="X", label="Y")
        # 첫 번째 Rectangle 의 위치
        rects = [sh for sh in s.shapes if "Rectangle" in sh.name]
        assert rects
        y_in = rects[0].top / 914400
        x_in = rects[0].left / 914400
        assert y_in > 2.0, f"METRIC_CARD y 변환 실패: y_in={y_in}"
        assert x_in > 0.5, f"METRIC_CARD x 변환 실패: x_in={x_in}"

    def test_highlight_with_float_y(self):
        sk, _, s = self._new_slide()
        sk.HIGHLIGHT(s, text="t", y=3.0)
        rounded = [sh for sh in s.shapes if "Rounded" in sh.name]
        assert rounded
        y_in = rounded[0].top / 914400
        assert y_in > 2.5, f"HIGHLIGHT y 변환 실패: y_in={y_in}"

    def test_quote_with_float_y(self):
        sk, _, s = self._new_slide()
        sk.QUOTE(s, text="t", y=4.0)
        # 첫 shape 의 y
        first = list(s.shapes)[0]
        y_in = first.top / 914400
        assert y_in > 3.5, f"QUOTE y 변환 실패: y_in={y_in}"

    def test_text_with_float_coords(self):
        sk, _, s = self._new_slide()
        sk.T(s, 0.6, 2.5, 7.07, 1.5, "test")
        tb = list(s.shapes)[0]
        x_in = tb.left / 914400
        y_in = tb.top / 914400
        w_in = tb.width / 914400
        assert abs(x_in - 0.6) < 0.01
        assert abs(y_in - 2.5) < 0.01
        assert abs(w_in - 7.07) < 0.01

    def test_a4_portrait_no_oob(self):
        """A4 세로에서 KPIS/HIGHLIGHT 등이 슬라이드 경계 내."""
        sk, prs, s = self._new_slide("delivery_a4_portrait")
        sk.KPIS(s, items=[{"value": "1", "label": "a"},
                          {"value": "2", "label": "b"},
                          {"value": "3", "label": "c"}], y=2.5)
        sk.HIGHLIGHT(s, text="test", y=6.0)
        sk.QUOTE(s, text="q", y=8.0)

        sw = prs.slide_width / 914400
        sh = prs.slide_height / 914400
        for shape in s.shapes:
            if shape.left is None:
                continue
            x = shape.left / 914400
            y = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            assert x >= -0.05, f"{shape.name} x={x}"
            assert y >= -0.05, f"{shape.name} y={y}"
            assert x + w <= sw + 0.05, f"{shape.name} x+w={x + w} > {sw}"
            assert y + h <= sh + 0.05, f"{shape.name} y+h={y + h} > {sh}"

        sk.apply_format("legacy_16_9")  # 복원


class TestSrcDynamicWidth:
    """SRC 가 hardcoded 8.0\" 대신 CW 사용."""

    def test_src_uses_cw(self):
        from src.generators import slide_kit
        slide_kit.apply_format("delivery_a4_portrait")
        prs = slide_kit.new_presentation()
        s = slide_kit.new_slide(prs)
        slide_kit.SRC(s, "test source")
        # SRC 가 만든 textbox
        tbs = [sh for sh in s.shapes if sh.has_text_frame
               and "Source:" in sh.text_frame.text]
        assert tbs
        w_in = tbs[0].width / 914400
        # A4 세로 CW = 8.27 - 0.6*2 = 7.07
        assert abs(w_in - 7.07) < 0.01, f"SRC width 동적 적용 실패: w_in={w_in}"
        slide_kit.apply_format("legacy_16_9")  # 복원


# ════════════════════════════════════════════════════════════
# THEMES — 10종 검증
# ════════════════════════════════════════════════════════════

class TestThemes:
    """10 테마 정의 + 일관성."""

    def test_ten_themes_defined(self):
        from src.generators.slide_kit import THEMES
        assert len(THEMES) == 10

    def test_legacy_six_themes_present(self):
        from src.generators.slide_kit import THEMES
        for name in [
            "warm_minimal", "classic_blue", "forest",
            "corporate", "mono_black", "soft_purple",
        ]:
            assert name in THEMES, f"기존 테마 {name} 누락"

    def test_new_four_themes_present(self):
        from src.generators.slide_kit import THEMES
        for name in [
            "sunset_orange", "ocean_teal", "gold_luxury", "arctic_white",
        ]:
            assert name in THEMES, f"신규 테마 {name} 누락"

    def test_each_theme_has_required_keys(self):
        """모든 테마는 14개 컬러 키 + _meta 보유."""
        from src.generators.slide_kit import THEMES
        required_color_keys = {
            "primary", "secondary", "teal", "accent", "dark", "light",
            "white", "gray", "lgray", "green", "orange", "gold",
            "card_bg", "card_border",
        }
        for name, theme in THEMES.items():
            missing = required_color_keys - set(theme.keys())
            assert not missing, f"{name}: 누락 키 {missing}"
            assert "_meta" in theme, f"{name}: _meta 누락"

    def test_each_theme_meta_complete(self):
        """모든 테마 _meta 는 margin/font_hero/font_action/font_body 보유."""
        from src.generators.slide_kit import THEMES
        for name, theme in THEMES.items():
            meta = theme["_meta"]
            for k in ["margin", "font_hero", "font_action", "font_body"]:
                assert k in meta, f"{name}._meta.{k} 누락"

    def test_color_values_are_rgb_tuples(self):
        """컬러 값은 (R,G,B) 0~255 tuple."""
        from src.generators.slide_kit import THEMES
        for name, theme in THEMES.items():
            for key, value in theme.items():
                if key.startswith("_"):
                    continue
                assert isinstance(value, tuple) and len(value) == 3, f"{name}.{key} != RGB tuple"
                for c in value:
                    assert 0 <= c <= 255, f"{name}.{key}={value} 범위 초과"


class TestApplyTheme:
    """``apply_theme()`` — 신규 테마 포함 동작."""

    @pytest.mark.parametrize("theme_name", [
        "warm_minimal", "classic_blue", "forest", "corporate",
        "mono_black", "soft_purple",
        "sunset_orange", "ocean_teal", "gold_luxury", "arctic_white",
    ])
    def test_apply_each_theme(self, theme_name):
        from src.generators import slide_kit
        result = slide_kit.apply_theme(theme_name)
        assert result == theme_name
        # 적용 후 C 가 새 컬러로 갱신
        assert slide_kit.C["primary"] is not None
        # warm_minimal 로 복원
        slide_kit.apply_theme("warm_minimal")

    def test_unknown_theme_raises(self):
        from src.generators import slide_kit
        with pytest.raises(ValueError, match="Unknown theme"):
            slide_kit.apply_theme("not_a_real_theme")


# ════════════════════════════════════════════════════════════
# LAYOUTS — 36종 검증
# ════════════════════════════════════════════════════════════

class TestNewLayouts:
    """신규 6 레이아웃 정의 + 동작."""

    @pytest.mark.parametrize("name", [
        "VERTICAL_STEPS", "TIMELINE_VERTICAL", "HEADLINE_NUMBER",
        "CARD_GRID_6", "AGENDA_TWO_COL", "PRICING_TABLE",
    ])
    def test_new_layout_in_LAYOUTS(self, name):
        from src.generators.slide_kit import LAYOUTS
        assert name in LAYOUTS

    def test_total_count_36(self):
        from src.generators.slide_kit import LAYOUTS
        assert len(LAYOUTS) == 36

    def test_vertical_steps_has_5_zones(self):
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["VERTICAL_STEPS"]["zones"]
        assert len(zones) == 5

    def test_timeline_vertical_has_8_zones(self):
        """4 단계 × 2 (date+event)"""
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["TIMELINE_VERTICAL"]["zones"]
        assert len(zones) == 8

    def test_headline_number_has_3_zones(self):
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["HEADLINE_NUMBER"]["zones"]
        assert len(zones) == 3
        ids = [z["id"] for z in zones]
        assert "big_number" in ids and "headline" in ids and "detail" in ids

    def test_card_grid_6_has_6_zones(self):
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["CARD_GRID_6"]["zones"]
        assert len(zones) == 6

    def test_agenda_two_col_has_10_zones(self):
        """좌5 + 우5 = 10 항목."""
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["AGENDA_TWO_COL"]["zones"]
        assert len(zones) == 10

    def test_pricing_table_has_9_zones(self):
        """3 플랜 × 3 (header + price + features)."""
        from src.generators.slide_kit import LAYOUTS
        zones = LAYOUTS["PRICING_TABLE"]["zones"]
        assert len(zones) == 9


class TestNewLayoutsAcrossFormats:
    """신규 6 layout 도 3 포맷 모두에서 경계 내."""

    @pytest.mark.parametrize("format_name, sw, sh", [
        ("legacy_16_9", 13.333, 7.5),
        ("delivery_a4_portrait", 8.27, 11.69),
        ("presentation_a4_landscape", 11.69, 8.27),
    ])
    def test_new_layouts_within_bounds(self, format_name, sw, sh):
        from src.generators import slide_kit
        slide_kit.apply_format(format_name)

        new_layouts = [
            "VERTICAL_STEPS", "TIMELINE_VERTICAL", "HEADLINE_NUMBER",
            "CARD_GRID_6", "AGENDA_TWO_COL", "PRICING_TABLE",
        ]
        for name in new_layouts:
            layout = slide_kit.LAYOUTS[name]
            for zone in layout["zones"]:
                x, y, w, h = zone["x"], zone["y"], zone["w"], zone["h"]
                assert x >= -0.01, f"{format_name}.{name}.{zone['id']}: x={x} < 0"
                assert y >= -0.01, f"{format_name}.{name}.{zone['id']}: y={y} < 0"
                assert x + w <= sw + 0.01, (
                    f"{format_name}.{name}.{zone['id']}: x+w={x + w:.2f} > {sw}"
                )
                assert y + h <= sh + 0.01, (
                    f"{format_name}.{name}.{zone['id']}: y+h={y + h:.2f} > {sh}"
                )

        slide_kit.apply_format("legacy_16_9")  # 복원
