"""
변환 에이전트 (Stage 5)

PPTX → JSON → Gemini + 디자인 레퍼런스 → HTML → Figma
"""

import json
import os
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from .base_agent import BaseAgent
from ..schemas.conversion_schema import (
    ConversionResult,
    HTMLSlide,
    ShapeData,
    SlideJSON,
)
from ..utils.logger import get_logger

logger = get_logger("conversion_agent")


class ConversionAgent(BaseAgent):
    """변환 에이전트 — PPTX → JSON → HTML → Figma"""

    def __init__(
        self,
        api_key: Optional[str] = None,
        model: Optional[str] = None,
        gemini_api_key: Optional[str] = None,
    ):
        super().__init__(api_key=api_key, model=model)
        self.gemini_api_key = gemini_api_key or os.getenv("GEMINI_API_KEY", "")

    async def execute(
        self,
        input_data: Dict[str, Any],
        progress_callback: Optional[Callable] = None,
    ) -> ConversionResult:
        """
        PPTX → JSON → HTML → Figma 변환

        Args:
            input_data: {
                "pptx_path": str,
                "output_dir": str,
                "design_reference": str (optional, PPTX/이미지 경로),
                "design_note": str (optional),
                "figma_import": bool (optional, default: False),
            }
        """
        pptx_path = input_data.get("pptx_path", "")
        output_dir = input_data.get("output_dir", "output/conversion")
        design_reference = input_data.get("design_reference")
        design_note = input_data.get("design_note", "")
        figma_import = input_data.get("figma_import", False)

        os.makedirs(output_dir, exist_ok=True)
        errors = []

        # ── Step 5a: PPTX → JSON ──
        if progress_callback:
            progress_callback({
                "stage": "conversion",
                "step": 1,
                "total": 3,
                "message": "PPTX → JSON 변환 중...",
            })

        slides_json = self._pptx_to_json(pptx_path)
        json_path = os.path.join(output_dir, "slides.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(
                [s.model_dump() for s in slides_json],
                f,
                ensure_ascii=False,
                indent=2,
            )

        logger.info(f"PPTX → JSON 완료: {len(slides_json)}개 슬라이드")

        # ── Step 5b: JSON → Gemini → HTML ──
        if progress_callback:
            progress_callback({
                "stage": "conversion",
                "step": 2,
                "total": 3,
                "message": "Gemini HTML 생성 중...",
            })

        html_slides = []
        try:
            html_slides = await self._json_to_html(
                slides_json, design_reference, design_note
            )
            # HTML 파일 저장
            html_dir = os.path.join(output_dir, "html")
            os.makedirs(html_dir, exist_ok=True)
            for slide in html_slides:
                html_path = os.path.join(html_dir, f"slide_{slide.slide_index:03d}.html")
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(slide.inline_html or slide.html_content)

            logger.info(f"HTML 생성 완료: {len(html_slides)}개 슬라이드")
        except Exception as e:
            errors.append(f"HTML 생성 실패: {str(e)}")
            logger.error(f"HTML 생성 실패: {e}")

        # ── Step 5c: Figma 준비 ──
        figma_url = None
        figma_ready = len(html_slides) > 0

        if progress_callback:
            progress_callback({
                "stage": "conversion",
                "step": 3,
                "total": 3,
                "message": "변환 완료",
            })

        return ConversionResult(
            json_data=slides_json,
            json_output_path=json_path,
            html_slides=html_slides,
            html_output_dir=os.path.join(output_dir, "html") if html_slides else None,
            figma_ready=figma_ready,
            figma_import_method="both",
            figma_url=figma_url,
            total_slides=len(slides_json),
            success=len(errors) == 0,
            errors=errors,
        )

    def _pptx_to_json(self, pptx_path: str) -> List[SlideJSON]:
        """PPTX 파일을 SlideJSON 리스트로 변환"""
        from pptx import Presentation
        from pptx.util import Inches, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slides_json = []
        prs = Presentation(pptx_path)

        slide_width = prs.slide_width / 914400  # EMU to inches
        slide_height = prs.slide_height / 914400

        for idx, slide in enumerate(prs.slides):
            shapes_data = []

            for shape in slide.shapes:
                shape_data = ShapeData(
                    shape_id=shape.shape_id,
                    shape_type=self._get_shape_type(shape),
                    left=shape.left / 914400 if shape.left else 0,
                    top=shape.top / 914400 if shape.top else 0,
                    width=shape.width / 914400 if shape.width else 0,
                    height=shape.height / 914400 if shape.height else 0,
                    rotation=shape.rotation if hasattr(shape, "rotation") else 0,
                )

                # 텍스트 추출
                if shape.has_text_frame:
                    shape_data.text = shape.text_frame.text
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                shape_data.font_name = run.font.name
                            if run.font.size:
                                shape_data.font_size = run.font.size / 12700  # EMU to pt
                            if run.font.bold:
                                shape_data.bold = True
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    shape_data.font_color = f"#{run.font.color.rgb}"
                            except (AttributeError, TypeError):
                                pass
                            break
                        break

                # 채우기 색상
                if hasattr(shape, "fill"):
                    try:
                        fill = shape.fill
                        if fill.type is not None:
                            shape_data.fill_type = str(fill.type)
                            if hasattr(fill, "fore_color") and fill.fore_color and fill.fore_color.rgb:
                                shape_data.fill_color = f"#{fill.fore_color.rgb}"
                    except Exception:
                        pass

                # 테이블
                if shape.has_table:
                    shape_data.shape_type = "table"
                    table = shape.table
                    table_data = {
                        "rows": table.rows.__len__(),
                        "cols": table.columns.__len__(),
                        "cells": [],
                    }
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            row_cells.append(cell.text)
                        table_data["cells"].append(row_cells)
                    shape_data.table_data = table_data

                shapes_data.append(shape_data)

            # 배경 추출
            background = {}
            try:
                bg = slide.background
                if bg.fill and bg.fill.type is not None:
                    background["type"] = str(bg.fill.type)
                    if hasattr(bg.fill, "fore_color") and bg.fill.fore_color.rgb:
                        background["color"] = f"#{bg.fill.fore_color.rgb}"
            except Exception:
                pass

            # 노트
            notes = None
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text

            slides_json.append(SlideJSON(
                slide_index=idx,
                width=slide_width,
                height=slide_height,
                shapes=shapes_data,
                background=background,
                notes=notes,
            ))

        return slides_json

    def _get_shape_type(self, shape) -> str:
        """도형 유형 문자열 변환"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        type_map = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
            MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
            MSO_SHAPE_TYPE.PICTURE: "image",
            MSO_SHAPE_TYPE.TABLE: "table",
            MSO_SHAPE_TYPE.CHART: "chart",
            MSO_SHAPE_TYPE.GROUP: "group",
            MSO_SHAPE_TYPE.FREEFORM: "freeform",
        }
        return type_map.get(shape.shape_type, "unknown")

    async def _json_to_html(
        self,
        slides_json: List[SlideJSON],
        design_reference: Optional[str],
        design_note: str,
    ) -> List[HTMLSlide]:
        """SlideJSON → Gemini → HTML 변환"""
        try:
            import google.generativeai as genai
        except ImportError:
            raise ImportError("google-generativeai 패키지가 필요합니다: pip install google-generativeai")

        if not self.gemini_api_key:
            raise ValueError("GEMINI_API_KEY가 설정되지 않았습니다")

        genai.configure(api_key=self.gemini_api_key)
        model = genai.GenerativeModel("gemini-2.0-flash")

        # 디자인 레퍼런스 분석
        reference_info = ""
        if design_reference:
            reference_info = self._analyze_reference(design_reference)

        html_slides = []

        # 배치 처리 (5개씩)
        batch_size = 5
        for batch_start in range(0, len(slides_json), batch_size):
            batch = slides_json[batch_start:batch_start + batch_size]
            batch_data = [s.model_dump() for s in batch]

            prompt = f"""
다음 프레젠테이션 슬라이드 JSON 데이터를 HTML+CSS로 변환해주세요.

## 슬라이드 데이터
{json.dumps(batch_data, ensure_ascii=False, indent=2)[:8000]}

## 디자인 지침
{design_note}
{reference_info}

## 변환 규칙
1. 각 슬라이드를 독립적인 HTML 파일로 생성
2. 슬라이드 크기: 1920px × 1080px (16:9)
3. 원본의 위치/크기/색상/폰트를 최대한 유지
4. CSS는 인라인 스타일로 적용
5. 디자인 레퍼런스가 있으면 시각적 품질 향상

JSON 배열로 응답해주세요:
```json
[
    {{
        "slide_index": 0,
        "html_content": "<div>...</div>",
        "css_content": "",
        "inline_html": "<!DOCTYPE html><html>...</html>"
    }}
]
```
"""

            try:
                response = model.generate_content(
                    prompt,
                    generation_config=genai.GenerationConfig(
                        temperature=0.2,
                        max_output_tokens=8192,
                    ),
                )

                # JSON 추출
                result_text = response.text
                data = self._extract_json(result_text)

                if isinstance(data, list):
                    for item in data:
                        html_slides.append(HTMLSlide(**item))
                elif isinstance(data, dict) and "slides" in data:
                    for item in data["slides"]:
                        html_slides.append(HTMLSlide(**item))

            except Exception as e:
                logger.error(f"Gemini HTML 생성 실패 (batch {batch_start}): {e}")

        return html_slides

    def _analyze_reference(self, reference_path: str) -> str:
        """디자인 레퍼런스 분석"""
        try:
            from ..utils.reference_analyzer import ReferenceAnalyzer
            analyzer = ReferenceAnalyzer()
            result = analyzer.analyze(reference_path)
            profile = analyzer.to_design_profile()
            return f"""
## 디자인 레퍼런스 분석 결과
- 주요 색상: {json.dumps(profile.get('theme', {}), ensure_ascii=False)}
- 폰트: {profile.get('primary_font', '미확인')}
- 스타일: {profile.get('style', '미확인')}
"""
        except Exception as e:
            logger.warning(f"레퍼런스 분석 실패: {e}")
            return ""
