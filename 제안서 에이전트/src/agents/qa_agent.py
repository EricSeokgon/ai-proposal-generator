"""
검수 에이전트 (Stage 4-2)

생성된 PPTX의 품질 검증
"""

import os
from typing import Any, Callable, Dict, List, Optional

from .base_agent import BaseAgent
from ..schemas.production_schema import QAReport, QAIssue
from ..utils.logger import get_logger

logger = get_logger("qa_agent")


class QAAgent(BaseAgent):
    """PPTX 검수 에이전트"""

    async def execute(
        self,
        input_data: Dict[str, Any],
        progress_callback: Optional[Callable] = None,
    ) -> QAReport:
        """
        생성된 PPTX 품질 검수

        Args:
            input_data: {
                "pptx_path": str,
                "plan": dict (ProposalPlan),
            }
        """
        pptx_path = input_data.get("pptx_path", "")
        plan = input_data.get("plan", {})

        if progress_callback:
            progress_callback({
                "stage": "qa",
                "step": 1,
                "total": 3,
                "message": "PPTX 파일 분석 중...",
            })

        issues: List[QAIssue] = []

        # Step 1: 파일 존재 확인
        if not os.path.exists(pptx_path):
            return QAReport(
                total_issues=1,
                critical_count=1,
                issues=[QAIssue(
                    slide_index=0,
                    severity="critical",
                    category="structure",
                    description=f"PPTX 파일이 존재하지 않습니다: {pptx_path}",
                    suggestion="ProductionAgent 재실행 필요",
                )],
                passed=False,
                summary="PPTX 파일 미생성",
            )

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu

            prs = Presentation(pptx_path)
        except Exception as e:
            return QAReport(
                total_issues=1,
                critical_count=1,
                issues=[QAIssue(
                    slide_index=0,
                    severity="critical",
                    category="structure",
                    description=f"PPTX 파일 열기 실패: {str(e)}",
                    suggestion="파일 손상 가능성. ProductionAgent 재실행",
                )],
                passed=False,
                summary="PPTX 파일 손상",
            )

        if progress_callback:
            progress_callback({
                "stage": "qa",
                "step": 2,
                "total": 3,
                "message": "품질 검증 수행 중...",
            })

        # Step 2: 검증 수행
        issues.extend(self._check_slide_count(prs, plan))
        issues.extend(self._check_fonts(prs))
        issues.extend(self._check_overlaps(prs))
        issues.extend(self._check_placeholders(prs))
        issues.extend(self._check_page_numbers(prs))
        issues.extend(self._check_empty_slides(prs))

        if progress_callback:
            progress_callback({
                "stage": "qa",
                "step": 3,
                "total": 3,
                "message": "검수 보고서 작성 중...",
            })

        # Step 3: 보고서 생성
        critical = [i for i in issues if i.severity == "critical"]
        warnings = [i for i in issues if i.severity == "warning"]
        infos = [i for i in issues if i.severity == "info"]

        passed = len(critical) == 0

        report = QAReport(
            total_issues=len(issues),
            critical_count=len(critical),
            warning_count=len(warnings),
            info_count=len(infos),
            issues=issues,
            passed=passed,
            summary=self._generate_summary(prs, critical, warnings),
            slide_count_match=not any(i.category == "structure" for i in critical),
            font_consistent=not any(i.category == "font" for i in issues),
            color_consistent=not any(i.category == "color" for i in issues),
            placeholder_format=not any(i.category == "placeholder" for i in issues),
            page_numbers_present=not any(i.category == "page_number" for i in issues),
        )

        logger.info(
            f"검수 완료: {'통과' if passed else '미통과'} "
            f"(critical={len(critical)}, warning={len(warnings)}, info={len(infos)})"
        )

        return report

    def _check_slide_count(self, prs, plan: Dict) -> List[QAIssue]:
        """슬라이드 수 검증"""
        issues = []
        expected = plan.get("structure", {}).get("total_slides", 0)
        actual = len(prs.slides)

        if expected > 0 and abs(actual - expected) > expected * 0.3:
            issues.append(QAIssue(
                slide_index=0,
                severity="critical",
                category="structure",
                description=f"슬라이드 수 불일치: 예상 {expected}장, 실제 {actual}장 (30% 이상 차이)",
                suggestion=f"기획된 {expected}장에 맞게 재생성",
            ))
        elif expected > 0 and abs(actual - expected) > expected * 0.1:
            issues.append(QAIssue(
                slide_index=0,
                severity="warning",
                category="structure",
                description=f"슬라이드 수 차이: 예상 {expected}장, 실제 {actual}장",
                suggestion="일부 슬라이드 누락 또는 추가 확인",
            ))

        return issues

    def _check_fonts(self, prs) -> List[QAIssue]:
        """폰트 일관성 검증"""
        issues = []
        non_standard_fonts = set()

        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_name = run.font.name
                        if font_name and font_name not in ("Pretendard", "Arial", None):
                            non_standard_fonts.add((idx, font_name))

        for slide_idx, font_name in list(non_standard_fonts)[:5]:
            issues.append(QAIssue(
                slide_index=slide_idx,
                severity="warning",
                category="font",
                description=f"비표준 폰트 사용: '{font_name}'",
                suggestion="Pretendard로 변경",
            ))

        return issues

    def _check_overlaps(self, prs) -> List[QAIssue]:
        """요소 겹침 검증"""
        issues = []

        for idx, slide in enumerate(prs.slides):
            shapes = []
            for shape in slide.shapes:
                shapes.append({
                    "left": shape.left,
                    "top": shape.top,
                    "right": shape.left + shape.width,
                    "bottom": shape.top + shape.height,
                    "name": shape.name,
                })

            for i in range(len(shapes)):
                for j in range(i + 1, len(shapes)):
                    a, b = shapes[i], shapes[j]
                    overlap_x = max(0, min(a["right"], b["right"]) - max(a["left"], b["left"]))
                    overlap_y = max(0, min(a["bottom"], b["bottom"]) - max(a["top"], b["top"]))

                    if overlap_x > 0 and overlap_y > 0:
                        overlap_area = overlap_x * overlap_y
                        a_area = (a["right"] - a["left"]) * (a["bottom"] - a["top"])
                        if a_area > 0 and overlap_area / a_area > 0.3:
                            issues.append(QAIssue(
                                slide_index=idx,
                                severity="warning",
                                category="overlap",
                                description=f"요소 겹침: '{a['name']}' ↔ '{b['name']}'",
                                suggestion="요소 위치 조정 필요",
                            ))
                            break
                if len(issues) > 20:
                    break

        return issues[:10]

    def _check_placeholders(self, prs) -> List[QAIssue]:
        """플레이스홀더 형식 검증"""
        issues = []
        import re

        bad_patterns = [
            (re.compile(r"OOO|XXX|___"), "OOO/XXX/___ 형식"),
            (re.compile(r"000-0000"), "전화번호 더미 데이터"),
        ]

        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text
                for pattern, desc in bad_patterns:
                    if pattern.search(text):
                        issues.append(QAIssue(
                            slide_index=idx,
                            severity="info",
                            category="placeholder",
                            description=f"비표준 플레이스홀더: {desc}",
                            suggestion="[대괄호] 형식으로 변경",
                        ))
                        break

        return issues[:10]

    def _check_page_numbers(self, prs) -> List[QAIssue]:
        """페이지 번호 존재 확인"""
        issues = []
        slides_without_pn = 0

        for idx, slide in enumerate(prs.slides):
            if idx < 2:  # 표지/목차는 제외
                continue
            has_page_number = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text.isdigit() and int(text) > 0:
                        has_page_number = True
                        break
            if not has_page_number:
                slides_without_pn += 1

        if slides_without_pn > len(prs.slides) * 0.5:
            issues.append(QAIssue(
                slide_index=0,
                severity="info",
                category="page_number",
                description=f"페이지 번호 미표기 슬라이드: {slides_without_pn}장",
                suggestion="PN() 함수로 페이지 번호 추가",
            ))

        return issues

    def _check_empty_slides(self, prs) -> List[QAIssue]:
        """빈 슬라이드 검증"""
        issues = []

        for idx, slide in enumerate(prs.slides):
            has_content = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    has_content = True
                    break
            if not has_content:
                issues.append(QAIssue(
                    slide_index=idx,
                    severity="warning",
                    category="content",
                    description=f"슬라이드 {idx+1}: 텍스트 콘텐츠 없음",
                    suggestion="콘텐츠 추가 또는 의도적 빈 슬라이드인지 확인",
                ))

        return issues[:5]

    def _generate_summary(self, prs, critical, warnings) -> str:
        """검수 요약 생성"""
        total = len(prs.slides)
        if not critical and not warnings:
            return f"검수 통과: {total}장 슬라이드, 이슈 없음"
        return (
            f"총 {total}장 슬라이드 검수 완료. "
            f"Critical {len(critical)}건, Warning {len(warnings)}건"
        )
